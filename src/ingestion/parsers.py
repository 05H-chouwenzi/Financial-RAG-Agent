"""文档解析器 —— 把 PDF/TXT/MD/DOCX/XLSX/PPTX/图片 解析成 LayoutBlock 列表

分层设计（对应规划 M2）：
1. 基线解析（light）：pdfplumber 抽取文本 + 表格（含二维结构）
2. 版面增强（heavy，可选）：MinerU / PaddleOCR PP-Structure —— 懒加载，装了才用
3. OCR：扫描页图片识别（src.ingestion.ocr）

统一产出 LayoutBlock，供结构感知切块使用。
"""
from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Optional

from src.ingestion.layout import BlockType, LayoutBlock
from src.ingestion.ocr import ocr_image, ocr_pdf_pages

logger = logging.getLogger("parsers")

# ============================================================
# 文本启发式
# ============================================================

_HEADING_RE = re.compile(
    r"^(第[一二三四五六七八九十百0-9]+[章节部篇]|"
    r"附录|目录|摘要|"
    r"[0-9]+(\.[0-9]+)*[、.．\s]|"
    r"[一二三四五六七八九十]+[、.．]|"
    r"\([一二三四五六七八九十]+\)|（[一二三四五六七八九十]+）)"
)

_ENDING_PUNCT = "。！？；，,.;:："


def _clean(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def _looks_like_heading(line: str) -> bool:
    """标题启发式：短行 + 标题模式 + 不以句号结尾"""
    line = line.strip()
    if len(line) == 0 or len(line) > 40:
        return False
    if line[-1] in _ENDING_PUNCT:
        return False
    if _HEADING_RE.match(line):
        return True
    # 文档/报告标题：短行且含 年度报告/半年报/公告 等字样
    if len(line) <= 40 and re.search(r"(年度报告|半年度报告|季度报告|公告|招股说明书|募集说明书)", line):
        return True
    # 全大写/数字开头的短行
    if re.match(r"^[A-Z0-9][A-Z0-9 \u4e00-\u9fff]{0,20}$", line):
        return True
    return False


def _heading_level(line: str) -> int:
    m = re.match(r"^第[一二三四五六七八九十百0-9]+[章节部篇]", line)
    if m:
        return 1
    m = re.match(r"^[0-9]+(\.[0-9]+)*", line)
    if m:
        return min(1 + line.count("."), 4)
    if re.match(r"^[一二三四五六七八九十]+[、.．]", line):
        return 2
    return 2


def _update_section(current: str, title: str, level: int) -> str:
    """维护章节路径（按层级覆盖）"""
    parts = [p for p in current.split("/") if p]
    parts = parts[: level - 1]
    parts.append(title.strip())
    return "/".join(parts)


def _make_block(
    block_type: str,
    text: str,
    page: int,
    order: int,
    section: str,
    doc_id: str,
    source: str,
    title_level: int = 0,
    table: Optional[list] = None,
) -> LayoutBlock:
    return LayoutBlock(
        doc_id=doc_id,
        block_type=block_type,
        text=text,
        table=table,
        page=page,
        reading_order=order,
        section_path=section,
        title_level=title_level,
        source=source,
    )


def _split_lines_to_blocks(
    lines: list[str], page: int, doc_id: str, source: str
) -> list[LayoutBlock]:
    """把文本行按标题/段落切分成 LayoutBlock"""
    blocks: list[LayoutBlock] = []
    order = 0
    buf: list[str] = []
    section = ""

    def flush():
        nonlocal buf
        if buf:
            blocks.append(_make_block(
                BlockType.PARAGRAPH.value, "\n".join(buf), page, order, section, doc_id, source
            ))
            buf = []

    for raw in lines:
        ln = _clean(raw)
        if not ln:
            continue
        if _looks_like_heading(ln):
            flush()
            level = _heading_level(ln)
            section = _update_section(section, ln, level)
            blocks.append(_make_block(
                BlockType.TITLE.value, ln, page, order, section, doc_id, source,
                title_level=level,
            ))
            order += 1
        else:
            buf.append(ln)
    flush()
    return blocks


def _merge_blocks(page_blocks: list[LayoutBlock]) -> list[LayoutBlock]:
    """合并同页相邻的同类型段落（减少碎片块）"""
    merged: list[LayoutBlock] = []
    for b in page_blocks:
        if b.block_type == BlockType.PARAGRAPH.value and merged and \
                merged[-1].block_type == BlockType.PARAGRAPH.value and \
                merged[-1].page == b.page and \
                merged[-1].section_path == b.section_path:
            merged[-1].text += "\n" + b.text
            continue
        merged.append(b)
    return merged


# ============================================================
# PDF 解析
# ============================================================

def _parse_pdf_light(pdf_path: Path, doc_id: str) -> list[LayoutBlock]:
    """基线解析：文本 + 表格（pdfplumber）"""
    import pdfplumber

    blocks: list[LayoutBlock] = []
    order = 0
    with pdfplumber.open(str(pdf_path)) as pdf:
        for page_idx, page in enumerate(pdf.pages, 1):
            # 表格（含二维结构）
            try:
                tables = page.extract_tables() or []
            except Exception as e:  # noqa: BLE001
                logger.warning("表格提取失败 page=%s: %s", page_idx, e)
                tables = []

            # 文本
            text = ""
            try:
                text = page.extract_text() or ""
            except Exception as e:  # noqa: BLE001
                logger.warning("文本提取失败 page=%s: %s", page_idx, e)

            # 页眉/页脚/页码过滤
            lines = [
                ln for ln in text.splitlines()
                if _clean(ln) and not re.fullmatch(r"[\d\s\-—/]{1,20}", _clean(ln))
            ]
            page_blocks = _split_lines_to_blocks(lines, page_idx, doc_id, str(pdf_path))
            for b in page_blocks:
                b.reading_order = order
                order += 1

            # 表格块追加（保留二维结构 + Markdown 渲染）
            for tbl in tables:
                rows = [[_clean(str(c)) if c is not None else "" for c in row] for row in tbl]
                if not any(any(c for c in row) for row in rows):
                    continue
                block = _make_block(
                    BlockType.TABLE.value, "", page_idx, order, "", doc_id, str(pdf_path),
                    table=rows,
                )
                block.text = block.table_to_markdown()
                block.reading_order = order
                blocks.append(block)
                order += 1

            blocks.extend(page_blocks)
    return _merge_blocks(blocks)


def _parse_pdf_heavy(pdf_path: Path, doc_id: str) -> Optional[list[LayoutBlock]]:
    """重武器解析（可选）：MinerU / PaddleOCR PP-Structure

    两个引擎都未安装时返回 None，自动回退基线解析。
    """
    # ---- PaddleOCR PP-StructureV3 ----
    try:
        from paddleocr import PPStructureV3  # type: ignore

        engine = PPStructureV3(show_log=False)
        result = engine.predict(str(pdf_path))
        blocks: list[LayoutBlock] = []
        order = 0
        for page in result:
            pno = int(getattr(page, "page", 0) or 0)
            for item in getattr(page, "res", None) or []:
                if not hasattr(item, "type"):
                    continue
                btype = item.type
                region = ""
                try:
                    region = item["res"]["text"] or ""
                except Exception:  # noqa: BLE001
                    region = ""
                if btype == "table":
                    html = item["res"]["html"] if "res" in item else ""
                    rows = _html_table_to_rows(html)
                    b = _make_block(
                        BlockType.TABLE.value, "", pno, order, "", doc_id, str(pdf_path), table=rows
                    )
                    b.text = b.table_to_markdown() if rows else html
                elif btype == "text" and region:
                    b = _make_block(BlockType.PARAGRAPH.value, region, pno, order, "", doc_id, str(pdf_path))
                else:
                    continue
                b.reading_order = order
                blocks.append(b)
                order += 1
        if blocks:
            logger.info("PaddleOCR PP-Structure 解析完成: %s", pdf_path.name)
            return _merge_blocks(blocks)
    except Exception as e:  # noqa: BLE001
        logger.debug("PP-Structure 不可用: %s", e)

    # ---- MinerU ----
    try:
        from magic_pdf.data.data_reader_writer import FileBasedDataWriter  # type: ignore
        from magic_pdf.pipe.UNIPipe import UNIPipe  # type: ignore

        # MinerU 2.x 接口，按需调整
        logger.warning("MinerU 集成点：请按你安装的版本对接 magic_pdf 输出")
    except Exception as e:  # noqa: BLE001
        logger.debug("MinerU 不可用: %s", e)

    return None


def _html_table_to_rows(html: str) -> list[list[str]]:
    """把 HTML 表格粗略转成二维列表（PP-Structure 输出）"""
    rows: list[list[str]] = []
    for tr in re.findall(r"<tr[^>]*>(.*?)</tr>", html, re.S):
        cells = re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", tr, re.S)
        rows.append([re.sub(r"<[^>]+>", "", c).strip() for c in cells])
    return rows


def parse_pdf(pdf_path: Path, doc_id: str, use_heavy: bool = False) -> list[LayoutBlock]:
    """解析 PDF：优先重武器，否则基线"""
    if use_heavy:
        heavy = _parse_pdf_heavy(pdf_path, doc_id)
        if heavy:
            return heavy
        logger.info("重武器不可用，回退基线解析: %s", pdf_path.name)
    return _parse_pdf_light(pdf_path, doc_id)


# ============================================================
# 其他格式
# ============================================================

def _parse_text_file(path: Path, doc_id: str) -> list[LayoutBlock]:
    raw = path.read_bytes()
    for enc in ("utf-8", "gbk", "gb18030", "utf-16"):
        try:
            text = raw.decode(enc)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    else:
        text = raw.decode("utf-8", errors="ignore")
    lines = [_clean(ln) for ln in text.splitlines()]
    blocks = _split_lines_to_blocks(lines, 0, doc_id, str(path))
    for i, b in enumerate(blocks):
        b.reading_order = i
    return blocks


def _parse_docx(path: Path, doc_id: str) -> list[LayoutBlock]:
    from docx import Document  # type: ignore

    doc = Document(str(path))
    blocks: list[LayoutBlock] = []
    order = 0
    section = ""
    for para in doc.paragraphs:
        text = _clean(para.text)
        if not text:
            continue
        style = (para.style.name or "").lower()
        if "heading" in style or _looks_like_heading(text):
            level = 1 if "heading 1" in style else 2
            section = _update_section(section, text, level)
            blocks.append(_make_block(BlockType.TITLE.value, text, 0, order, section, doc_id, str(path), title_level=level))
        else:
            blocks.append(_make_block(BlockType.PARAGRAPH.value, text, 0, order, section, doc_id, str(path)))
        order += 1
    # 表格
    for tbl in doc.tables:
        rows = [[_clean(c.text) for c in row.cells] for row in tbl.rows]
        b = _make_block(BlockType.TABLE.value, "", 0, order, section, doc_id, str(path), table=rows)
        b.text = b.table_to_markdown()
        blocks.append(b)
        order += 1
    return _merge_blocks(blocks)


def _parse_xlsx(path: Path, doc_id: str) -> list[LayoutBlock]:
    from openpyxl import load_workbook  # type: ignore

    wb = load_workbook(path, read_only=True, data_only=True)
    blocks: list[LayoutBlock] = []
    order = 0
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = [
            ["" if c is None else str(c) for c in row]
            for row in ws.iter_rows(values_only=True)
            if any(c is not None and str(c).strip() for c in row)
        ]
        if not rows:
            continue
        b = _make_block(BlockType.TABLE.value, "", 0, order, f"{sheet}", doc_id, str(path), table=rows)
        b.text = b.table_to_markdown()
        blocks.append(b)
        order += 1
    return blocks


def _parse_pptx(path: Path, doc_id: str) -> list[LayoutBlock]:
    from pptx import Presentation  # type: ignore

    prs = Presentation(str(path))
    blocks: list[LayoutBlock] = []
    order = 0
    for slide_no, slide in enumerate(prs.slides, 1):
        section = f"Slide {slide_no}"
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = _clean("\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip()))
                if text:
                    blocks.append(_make_block(BlockType.PARAGRAPH.value, text, slide_no, order, section, doc_id, str(path)))
                    order += 1
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = [[_clean(c.text) for c in row.cells] for row in shape.table.rows]
                b = _make_block(BlockType.TABLE.value, "", slide_no, order, section, doc_id, str(path), table=rows)
                b.text = b.table_to_markdown()
                blocks.append(b)
                order += 1
    return blocks


# ============================================================
# 统一入口
# ============================================================

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}

SUPPORTED_SUFFIXES = {".pdf", ".txt", ".md", ".markdown", ".docx", ".xlsx", ".xls", ".pptx"} | IMAGE_SUFFIXES


def parse_document(path, doc_id: str = "", use_heavy: bool = False) -> list[LayoutBlock]:
    """按扩展名解析文档，返回 LayoutBlock 列表"""
    path = Path(path)
    doc_id = doc_id or str(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(path, doc_id, use_heavy=use_heavy)
    if suffix in (".md", ".markdown", ".txt"):
        return _parse_text_file(path, doc_id)
    if suffix == ".docx":
        return _parse_docx(path, doc_id)
    if suffix in (".xlsx", ".xls"):
        return _parse_xlsx(path, doc_id)
    if suffix == ".pptx":
        return _parse_pptx(path, doc_id)
    if suffix in IMAGE_SUFFIXES:
        text = ocr_image(path)
        blocks = _split_lines_to_blocks(text.splitlines(), 0, doc_id, str(path)) if text else []
        for i, b in enumerate(blocks):
            b.reading_order = i
        return blocks
    raise ValueError(f"不支持的文件类型: {suffix}")
